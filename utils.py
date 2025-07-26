from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE

COLOR_PALETTE = [
    RGBColor(31, 119, 180),  # blue
    RGBColor(255, 127, 14),  # orange
    RGBColor(44, 160, 44),   # green
    RGBColor(214, 39, 40),   # red
    RGBColor(148, 103, 189), # purple
]

SUPPORTED_TYPES = {
    XL_CHART_TYPE.BAR_CLUSTERED,
    XL_CHART_TYPE.BAR_STACKED,
}

def safe_get_values(series):
    try:
        return [v for v in series.values if v is not None]
    except:
        return []

def process_pptx(path_in, seg_count, seg_names):
    prs = Presentation(path_in)
    color_map = COLOR_PALETTE[:seg_count]

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not shape.has_chart:
                continue

            chart = shape.chart
            if chart.chart_type not in SUPPORTED_TYPES:
                continue  # Skip unsupported chart types

            try:
                categories = [pt.label for pt in chart.plots[0].categories]
                original_series = chart.series
                if not original_series:
                    continue

                new_data = ChartData()
                new_data.categories = categories

                for original_series_index, series in enumerate(original_series):
                    original_name = series.name
                    original_values = safe_get_values(series)
                    if not original_values:
                        continue

                    # Create seg_count new series per original
                    for seg_index in range(seg_count):
                        seg_name = f"{original_name} - {seg_names[seg_index]}"
                        seg_values = original_values  # Use same values for all segments
                        new_data.add_series(seg_name, seg_values)

                chart.replace_data(new_data)

                # Set colors (each seg_index across series)
                new_series = chart.series
                for idx, series in enumerate(new_series):
                    color_idx = idx % seg_count  # So colors repeat per segment
                    series.name = new_data.series[idx][0]  # Set name
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = color_map[color_idx]

            except Exception as e:
                print(f"Error on slide {slide_index} shape {shape_index}: {e}")
                continue

    output_path = path_in.replace(".pptx", "_segmented.pptx")
    prs.save(output_path)
    return output_path
